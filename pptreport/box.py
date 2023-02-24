import os
import tempfile

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt

# For reading pictures
from PIL import Image
import fitz

# For fonts
import matplotlib.font_manager


def resize_text(txt_frame, max_size=18):
    """ Resize text to fit the textbox.

    Parameters
    ----------
    txt_frame : pptx.text.text.TextFrame
        The text frame to be resized.
    max_size : int, default 18
        The maximum fontsize of the text.
    """

    # Find all fonts and fit text to the box
    fonts = matplotlib.font_manager.findSystemFonts()

    # Find a font that works
    for font in fonts:
        try:
            txt_frame.fit_text(font_file=font, max_size=max_size)  # some fonts return a 'cannot unpack non-iterable NoneType object'-error
            break
        except Exception:
            pass


class Box():
    """ A box is a constrained area of the slide which contains a single element e.g. text, a picture, a table, etc. """

    def __init__(self, slide, coordinates):
        """
        Initialize a box.

        Parameters
        ----------
        slide : pptx slide object
            The slide on which the box is located.
        coordinates : tuple
            Coordinates containing (left, top, width, height) of the box (in pptx units).
        """

        self.slide = slide
        self.logger = None

        # Bounds of the box
        self.left = int(coordinates[0])
        self.top = int(coordinates[1])
        self.width = int(coordinates[2])
        self.height = int(coordinates[3])

        # Initialize bounds of the content (can be smaller than the box)
        self.content = None
        self.content_left = self.left
        self.content_top = self.top
        self.content_width = self.width
        self.content_height = self.height

        self.border = None  # border object of the box

    def add_parameters(self, parameters):
        """ Add parameters from the slide """

        for key, value in parameters.items():
            setattr(self, key, value)

    def add_border(self):
        """ Adds a border shape of box to make debugging easier """

        if self.border is None:
            self.border = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.left, self.top, self.width, self.height)
            self.border.fill.background()
            self.border.line.color.rgb = RGBColor(0, 0, 0)  # black
            self.border.line.width = Pt(1)

    def remove_border(self):
        """ Removes the border shape """

        if self.border is not None:
            self.border._sp.getparent().remove(self.border._sp)
            self.border = None  # reset border

    def fill(self, content, box_index=0):
        """
        Fill the box with content. The function estimates type of content is given.

        Parameters
        ----------
        content : str
            The element to be added to the box.
        box_index : int
            The index of the box (used for accessing properties per box).
        """

        self.box_index = box_index
        self.content = content

        # Find out what type of content it is
        content_type = self._get_content_type(content)

        if content_type == "pdf":
            filename = self.convert_pdf(content)
            self.fill_image(filename)
            os.remove(filename)

        elif content_type == "image":
            self.fill_image(content)

        elif content_type == "textfile":
            with open(content) as f:
                text = f.read()
            self.fill_text(text)

        elif content_type == "text":
            self.fill_text(content)

        elif content_type == "empty":
            return  # do nothing

        else:
            pass

        self.logger.debug(f"Box index {box_index} was filled with {content_type}")

    @staticmethod
    def _get_content_type(content):
        """ Determine the type of content. """

        if isinstance(content, str):
            if os.path.isfile(content):

                # Find out the content of the file
                try:
                    with open(content) as f:
                        _ = f.read()
                    return "textfile"

                except UnicodeDecodeError:

                    if content.endswith(".pdf"):
                        return "pdf"
                    else:
                        return "image"
            else:
                return "text"
        elif content is None:
            return "empty"
        else:
            t = type(content)
            raise ValueError(f"Content of type '{t}' cannot be added to slide")

    def convert_pdf(self, pdf):
        """ Convert a pdf file to a png file. """

        # Create temporary file
        temp_name = next(tempfile._get_candidate_names()) + ".png"
        temp_dir = tempfile.gettempdir()
        temp_file = os.path.join(temp_dir, temp_name)
        self.logger.debug(f"Converting pdf to temporary png at: {temp_file}")

        # Convert pdf to png
        doc = fitz.open(pdf)
        page = doc.load_page(0)
        pix = page.get_pixmap()
        pix.save(temp_file)

        return temp_file

    def fill_image(self, filename):
        """ Fill the box with an image. """

        # Find out the size of the image
        self._adjust_image_size(filename)
        self._adjust_image_position()  # adjust image position to middle of box

        # Add image
        self.logger.debug("Adding image to slide from file: " + filename)
        self.picture = self.slide.shapes.add_picture(filename, self.content_left, self.content_top, self.content_width, self.content_height)

    def _adjust_image_size(self, filename):
        """
        Adjust the size of the image to fit the box.

        Parameters
        ----------
        filename : str
            Path to the image file.
        """

        # Find out the size of the image
        im = Image.open(filename)
        im_width, im_height = im.size

        box_width = self.width
        box_height = self.height

        im_ratio = im_width / im_height  # >1 for landscape, <1 for portrait
        box_ratio = box_width / box_height

        # width is the limiting factor; height will be smaller than box_height
        if box_ratio < im_ratio:  # box is wider than image; height will be limiting
            self.content_width = box_width
            self.content_height = box_width * im_height / im_width  # maintain aspect ratio

        # height is the limiting factor; width will be smaller than box_width
        else:
            self.content_width = box_height * im_width / im_height  # maintain aspect ratio
            self.content_height = box_height

    def _get_content_alignment(self):
        """ Get the content alignment for this box. """

        self.logger.debug(f"Getting content alignment for box '{self.box_index}'. Input content alignment is '{self.content_alignment}'")

        if isinstance(self.content_alignment, str):  # if content alignment is a string, use it for all boxes
            this_alignment = self.content_alignment

        elif isinstance(self.content_alignment, list):  # if content alignment is a list, use the alignment for the current box
            if self.box_index > len(self.content_alignment) - 1:  # if box index is out of range, use default alignment
                this_alignment = "center"  # default alignment
            else:
                this_alignment = self.content_alignment[self.box_index]
        else:
            raise ValueError(f"Content alignment '{self.content_alignment}' is not valid. Valid content alignments are: str or list of str")

        # Check if current alignment is valid
        valid_alignments = ["left", "right", "center", "lower", "upper",
                            "lower left", "lower center", "lower right",
                            "upper left", "upper center", "upper right",
                            "center left", "center center", "center right"]

        if this_alignment.lower() not in valid_alignments:
            raise ValueError(f"Alignment '{self.content_alignment}' is not valid. Valid content alignments are: {valid_alignments}")

        # Expand into the structure "<vertical> <horizontal>"
        if this_alignment.lower() in ["left", "right", "center"]:
            this_alignment = "center " + this_alignment
        elif this_alignment.lower() in ["lower", "upper"]:
            this_alignment = this_alignment + " center"

        return this_alignment.split(" ")

    def _adjust_image_position(self):
        """ Adjust the position of the image to be in the middle of the box. """

        # Get content alignment for this box
        vertical, horizontal = self._get_content_alignment()

        # Adjust image position vertically
        if vertical == "upper":
            self.content_top = self.top
        elif vertical == "lower":
            self.content_top = self.top + self.height - self.content_height
        elif vertical == "center":
            self.content_top = self.top + (self.height - self.content_height) / 2

        # Adjust image position horizontally
        if horizontal == "left":
            self.content_left = self.left
        elif horizontal == "right":
            self.content_left = self.left + self.width - self.content_width
        elif horizontal == "center":
            self.content_left = self.left + (self.width - self.content_width) / 2

    def fill_text(self, text):
        """
        Fill the box with text.

        Parameters
        ----------
        text : str
            The text to be added to the box.
        """

        txt_box = self.slide.shapes.add_textbox(self.left, self.top, self.width, self.height)
        txt_frame = txt_box.text_frame
        txt_frame.word_wrap = True

        # Place all text in one paragraph
        # txt_frame.add_paragraph() # text_frame already has one paragraph
        p = txt_frame.paragraphs[0]
        p.text = text

        # Try to fit text size to the box
        resize_text(txt_frame)
        txt_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # An additional step of resizing text to fit the box

        # Set alignment of text in textbox
        vertical, horizontal = self._get_content_alignment()

        if vertical == "upper":
            txt_frame.vertical_anchor = MSO_ANCHOR.TOP
        elif vertical == "lower":
            txt_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
        elif vertical == "center":
            txt_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        if horizontal == "left":
            p.alignment = PP_ALIGN.LEFT
        elif horizontal == "right":
            p.alignment = PP_ALIGN.RIGHT
        elif horizontal == "center":
            p.alignment = PP_ALIGN.CENTER
