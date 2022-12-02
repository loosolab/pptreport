import os
import numpy as np
import glob
import tempfile
import inspect
import yaml

# For reading pictures
from PIL import Image
import fitz

# Pptx modules
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor


def flatten_list(lst):
    """ Flatten a list containing both lists and non-lists """

    flat = []
    for element in lst:
        if isinstance(element, list):
            flat.extend(element)
        else:
            flat.append(element)

    return flat


def get_default_args(func):
    signature = inspect.signature(func)
    defaults = {k: v.default for k, v in signature.parameters.items() if v.default is not inspect.Parameter.empty}

    return defaults

def glob_files(lst):

    content = [glob.glob(c) if "*" in c else c for c in lst]
    content = flatten_list(content)  # flatten list in case glob extended files

    return content


class PowerPointReport():

    def __init__(self, template=None, size="standard"):
        """ Initialize a presentation object using an existing presentation (template) or from scratch (default) """

        self.template = template
        self._prs = Presentation(template)
        self.size = size

        # Set size of the presentation
        if template is None:
            self.set_size(size)

        self._slides = []        # a list of SlidePlus objects
        self.borders = False    # show borders of content boxes

    def set_size(self, size):
        """
        Set the size of the presentation.

        Parameters
        ----------
        size : str or tuple of float
            Size of the presentation. Can be "standard", "widescreen", "a4-portait" or "a4-landscape". Can also be a tuple of numbers indicating (height, width) in cm.
        """

        if isinstance(size, tuple):
            if len(size) != 2:
                raise ValueError("Size tuple must be of length 2.")

            h, w = Cm(size[0]), Cm(size[1])

        elif size == "standard":
            h, w = Cm(19.05), Cm(25.4)

        elif size == "widescreen":
            h, w = Cm(19.05), Cm(33.867)

        elif size == "a4-portrait":
            h, w = Cm(27.517), Cm(19.05)

        elif size == "a4-landscape":
            h, w = Cm(19.05), Cm(27.517)

        else:
            raise ValueError("Invalid size given. Choose from: 'standard', 'widescreen', 'a4-portrait', 'a4-landscape' or a tuple of floats.")

        self._prs.slide_height = h
        self._prs.slide_width = w

    def add_title_slide(self, title, layout=0, subtitle=None):
        """
        Add a title slide to the presentation.

        Parameters
        ----------
        title : str
            Title of the slide.
        layout : int, default 0
            The layout of the slide. The first layout (0) is usually the default title slide.
        subtitle : str, optional
            Subtitle of the slide if the layout has a suptitle placeholder.
        """

        self.add_slide(title=title, slide_layout=layout)
        slide = self._slides[-1]

        # Fill placeholders
        if subtitle is not None:
            if len(slide.placeholders) == 2:
                slide.placeholders[1].text = subtitle

    def add_slide(self,
                  content=[],
                  title=None,
                  slide_layout=1,
                  content_layout="grid",
                  outer_margin=2,
                  inner_margin=1,
                  n_columns=2,
                  width_ratios=None,
                  height_ratios=None,
                  fontsize=12,
                  split=False,
                  ):
        """
        Add a slide to the presentation.

        Parameters
        ----------
        content : list of str
            List of content to be added to the slide. Can be either a path to a file or a string.
        title : str, optional
            Title of the slide.
        slide_layout : int, default 1
            Layout of the slide.
        content_layout : str, default "grid"
            Layout of the slide. Can be "grid", "vertical" or "horizontal". Can also be a numpy array of integers indicating the layout of the slide.
        outer_margin : float, default 2
            Outer margin of the slide (in cm).
        inner_margin : float, default 1
            Inner margin of the slide elements (in cm).
        n_columns : int, default 2
            Number of columns in the layout in case of "grid" layout.
        width_ratios : list of float, optional
            Width of the columns in case of "grid" layout.
        height_ratios : list of float, optional
            Height of the rows in case of "grid" layout.
        fontsize : int, default 12
            Fontsize of the text in the slide.
        split : bool, default False
            Split the content into multiple slides.
        """

        # Get all parameters
        parameters = locals()

        # Establish content
        if isinstance(content, str):
            content = [content]

        # If split is false, content should be contained in one slide
        if split is False:
            content = [content]
        else:
            content = glob_files(content)

        # Create slide(s)
        for slide_content in content:

            # Setup an empty slide
            slide = self.setup_slide(slide_layout)
            slide.add_parameters(parameters)
            slide.content = slide_content

            # Glob files
            slide._content = glob_files(slide_content)  # internal extension of content

            # Set title of slide
            slide.set_title(title)

            # Add content to slide
            if len(slide_content) > 0:

                slide.set_layout_matrix()  # Find the layout of the slide
                slide.create_boxes()       # Create boxes based on layout
                slide.fill_boxes()         # Fill boxes with content

    def setup_slide(self, slide_layout):
        """ Initialize an empty slide with a given layout. """

        layout_obj = self._prs.slide_layouts[slide_layout]
        slide_obj = self._prs.slides.add_slide(layout_obj)

        slide = Slide(slide_obj)
        slide._slide_height = self._prs.slide_height
        slide._slide_width = self._prs.slide_width
        self._slides.append(slide)

        return slide

    def show_borders(self):
        """ Show borders of all content boxes. Useful for debugging layouts."""

        for slide in self._slides:
            if hasattr(slide, "boxes"):
                for box in slide.boxes:
                    box.add_border()

        self.borders = True

    def get_config(self, full=False):
        """ Return a dictionary with the configuration of the presentation """

        config = {}

        for key in self.__dict__:
            if key[0] != "_":
                config[key] = self.__dict__[key]

        # Get default slide config
        defaults = get_default_args(self.add_slide)

        # Get config of each slide
        config["slides"] = []
        for slide in self._slides:
            slide_config = {}
            for key, value in slide.__dict__.items():
                if not key.startswith("_"):
                    slide_config[key] = slide.__dict__[key]

                    if full is False:
                        if value == defaults[key]:
                            del slide_config[key]
                        elif isinstance(value, list) and len(value) == 0:  # content can be an empty list
                            del slide_config[key]

            config["slides"].append(slide_config)

        return config

    def write_config(self, filename):

        config = self.get_config()

        with open(filename, "w") as f:
            yaml.dump(config, f)

    def from_config(self, config):
        """ Create a presentation from a configuration dictionary.

        Parameters
        ----------
        config : dict
            Configuration dictionary.
        """

        self.set_size(config["size"])

        for slide in config["slides"]:
            self.add_slide(title=slide["title"], slide_layout=slide["layout"])

    def save(self, filename):
        """
        Save the presentation to a file.

        Parameters
        ----------
        filename : str
            Filename of the presentation.
        """

        self._prs.save(filename)


# ------------------------------------------------------------------------------

class Slide():
    """ An internal class for creating slides. """

    def __init__(self, slide):

        self._slide = slide  # Slide object from python-pptx
        self._boxes = []   # Boxes in the slide

    def set_title(self, title):
        """ Set the title of the slide """

        if title is not None:
            self._slide.shapes.title.text = title

    def add_parameters(self, parameters):
        """ Add parameters to the slide """

        for key in parameters:
            if key != "self":
                setattr(self, key, parameters[key])

    def set_layout_matrix(self):
        """
        Get the content layout matrix for the slide.
        """

        # Get variables from self
        layout = self.content_layout
        n_elements = len(self._content)
        n_columns = self.n_columns

        # Get layout matrix depending on "layout" variable
        if layout == "grid":
            n_columns = min(n_columns, n_elements)  # number of columns cannot be larger than number of elements
            n_rows = int(np.ceil(n_elements / n_columns))  # number of rows to fit elements
            n_total = n_rows * n_columns

            intarray = list(range(n_elements))
            intarray.extend([np.nan] * (n_total - n_elements))

            layout_matrix = np.array(intarray).reshape((n_rows, n_columns))

        elif layout == "vertical":
            layout_matrix = np.array(list(range(n_elements))).reshape((n_elements, 1))

        elif layout == "horizontal":
            layout_matrix = np.array(list(range(n_elements))).reshape((1, n_elements))

        else:
            layout_matrix = self._validate_layout(layout)  # check if layout is a valid matrix

        self._layout_matrix = layout_matrix

    @staticmethod
    def _validate_layout(layout_matrix):
        """ Validate the given layout matrix. """
        # TODO: check if layout is a valid matrix

        layout_matrix = np.array(layout_matrix)

        return layout_matrix

    def create_boxes(self):
        """
        Create boxes for the slide dependent on the intrnal layout matrix.
        """

        layout_matrix = self._layout_matrix
        nrows, ncols = layout_matrix.shape

        # Convert margins from cm to pptx units
        outer_margin_unit = Cm(self.outer_margin)
        inner_margin_unit = Cm(self.inner_margin)
        left_margin = right_margin = outer_margin_unit
        top_margin = bottom_margin = outer_margin_unit

        # Add to top margin based on title
        if self._slide.shapes.title.text != "":
            top_margin = self._slide.shapes.title.top + self._slide.shapes.title.height + outer_margin_unit
        else:
            top_margin = outer_margin_unit
            sp = self._slide.shapes.title.element
            sp.getparent().remove(sp)

        # How many columns and rows are there?
        n_rows, n_cols = layout_matrix.shape

        # Get total height and width of pictures
        total_width = self._slide_width - left_margin - right_margin - (n_cols - 1) * inner_margin_unit
        total_height = self._slide_height - top_margin - bottom_margin - (n_rows - 1) * inner_margin_unit

        # Get column widths and row heights
        if self.width_ratios is None:
            widths = (np.ones(ncols) / ncols) * total_width
        else:
            widths = np.array(self.width_ratios) / sum(self.width_ratios) * total_width

        if self.height_ratios is None:
            heights = (np.ones(nrows) / nrows) * total_height
        else:
            heights = np.array(self.height_ratios) / sum(self.height_ratios) * total_height

        # Box coordinates
        box_numbers = layout_matrix[~np.isnan(layout_matrix)].flatten()
        box_numbers = sorted(set(box_numbers))  # unique box numbers
        box_numbers = [box_number for box_number in box_numbers if box_number >= 0]  # remove negative numbers
        for i in box_numbers:

            # Get column and row number
            coordinates = np.argwhere(layout_matrix == i)

            # Get upper left corner of box
            row, col = coordinates[0]
            left = left_margin + np.sum(widths[:col]) + col * inner_margin_unit
            top = top_margin + np.sum(heights[:row]) + row * inner_margin_unit

            # Get total width and height of box (can span multiple columns and rows)
            width = 0
            height = 0
            rows = set(coordinates[:, 0])
            for row in rows:
                height += heights[row]
            height += (len(rows) - 1) * inner_margin_unit  # add inner margins between rows

            cols = set(coordinates[:, 1])
            for col in cols:
                width += widths[col]
            width += (len(cols) - 1) * inner_margin_unit  # add inner margins between columns

            #  Create box
            self.add_box((left, top, width, height))

    def add_box(self, coordinates):
        """
        Add a box to the slide.

        Parameters
        ----------
        coordinates : tuple
            Coordinates containing (left, top, width, height) of the box (in pptx units).
        """

        box = Box(self._slide, coordinates)
        self._boxes.append(box)

    def fill_boxes(self):
        """ Fill the boxes with the elements in self.content """

        for i, element in enumerate(self._content):
            self._boxes[i].fill(element)


class Box():
    """ A box is a constrained area of the slide which contains a single element e.g. text, a picture, a table, etc. """

    def __init__(self, slide, coordinates):
        """
        Initialize a box.

        Parameters
        ----------
        slide : pptx slide object
            The slide on which the box is located.
        coordinates : tuple
            Coordinates containing (left, top, width, height) of the box (in pptx units)."""

        self.slide = slide

        # Bounds of the box
        self.left = int(coordinates[0])
        self.top = int(coordinates[1])
        self.width = int(coordinates[2])
        self.height = int(coordinates[3])

        # Initialize bounds of the content (can be smaller than the box)
        self.content = None
        self.content_left = self.left
        self.content_top = self.top
        self.content_width = self.width
        self.content_height = self.height

        self.border = None  # border object of the box

    def add_border(self):
        """ Adds a border shape of box to make debugging easier """

        if self.border is None:
            self.border = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.left, self.top, self.width, self.height)
            self.border.fill.background()
            self.border.line.color.rgb = RGBColor(0, 0, 0)
            self.border.line.width = Pt(1)

    def remove_border(self):
        """ Removes the border shape """

        if self.border is not None:
            self.border.getparent().remove(self.border)

    def fill(self, content):
        """
        Fill the box with content. The function estimates type of content is given.

        Parameters
        ----------
        content : str
            The element to be added to the box.
        """

        self.content = content

        # Find out what type of content it is
        content_type = self._get_content_type(content)

        if content_type == "pdf":
            filename = self.convert_pdf(content)
            self.fill_image(filename)
            os.remove(filename)

        elif content_type == "image":
            self.fill_image(content)

        elif content_type == "textfile":
            with open(content) as f:
                text = f.read()
            self.fill_text(text)

        elif content_type == "text":
            self.fill_text(content)

        else:
            pass

    @staticmethod
    def _get_content_type(content):
        """ Determine the type of content. """

        if isinstance(content, str):
            if os.path.isfile(content):

                # Find out the content of the file
                try:
                    with open(content) as f:
                        _ = f.read()
                    return "textfile"

                except UnicodeDecodeError:

                    if content.endswith(".pdf"):
                        return "pdf"
                    else:
                        return "image"
            else:
                return "text"
        else:
            t = type(content)
            raise ValueError(f"Content of type '{t}' cannot be added to slide")

    @staticmethod
    def convert_pdf(pdf):
        """ Convert a pdf file to a png file. """

        # Create temporary file
        temp_name = next(tempfile._get_candidate_names()) + ".png"
        temp_dir = tempfile.gettempdir()
        temp_file = os.path.join(temp_dir, temp_name)
        print(temp_file)

        # Convert pdf to png
        doc = fitz.open(pdf)
        page = doc.load_page(0)
        pix = page.get_pixmap()
        pix.save(temp_file)

        return temp_file

    def fill_image(self, filename):
        """ Fill the box with an image. """

        # Find out the size of the image
        self._adjust_image_size(filename)
        self._adjust_image_position()  # adjust image position to middle of box

        # Add image
        self.picture = self.slide.shapes.add_picture(filename, self.content_left, self.content_top, self.content_width, self.content_height)

    def _adjust_image_size(self, filename):
        """
        Adjust the size of the image to fit the box.

        Parameters
        ----------
        filename : str
            Path to the image file.
        """

        # Find out the size of the image
        im = Image.open(filename)
        im_width, im_height = im.size

        box_width = self.width
        box_height = self.height

        im_ratio = im_width / im_height  # >1 for landscape, <1 for portrait
        box_ratio = box_width / box_height

        # width is the limiting factor; height will be smaller than box_height
        if box_ratio < im_ratio:  # box is wider than image; height will be limiting
            self.content_width = box_width
            self.content_height = box_width * im_height / im_width  # maintain aspect ratio

        # height is the limiting factor; width will be smaller than box_width
        else:
            self.content_width = box_height * im_width / im_height  # maintain aspect ratio
            self.content_height = box_height

    def _adjust_image_position(self):
        """ Adjust the position of the image to be in the middle of the box. """

        self.content_left = self.left + (self.width - self.content_width) / 2
        self.content_top = self.top + (self.height - self.content_height) / 2

    def fill_text(self, text):
        """
        Fill the box with text.

        Parameters
        ----------
        text : str
            The text to be added to the box.
        """

        txt_box = self.slide.shapes.add_textbox(self.left, self.top, self.width, self.height)
        txt_frame = txt_box.text_frame

        txt_frame.add_paragraph()
        p = txt_frame.paragraphs[0]
        p.text = text
        txt_frame.word_wrap = True
        txt_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        txt_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        # txt_frame.fit_text() # only on windows?
