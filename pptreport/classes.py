import os
import numpy as np
import glob
import tempfile
import inspect
import pprint
import json
import re
import subprocess
import logging
import sys

# For reading pictures
from PIL import Image
import fitz

# For fonts
import matplotlib.font_manager

# Pptx modules
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor

###############################################################################
# ---------------------------- Helper functions ----------------------------- #
###############################################################################


def flatten_list(lst):
    """ Flatten a list containing both lists and non-lists """

    flat = []
    for element in lst:
        if isinstance(element, list):
            flat.extend(element)
        else:
            flat.append(element)

    return flat


def get_default_args(func):
    signature = inspect.signature(func)
    defaults = {k: v.default for k, v in signature.parameters.items() if v.default is not inspect.Parameter.empty}

    return defaults


def glob_files(lst):
    """ Glob files in a list of strings which might contain "*". If no files are found for glob, raise an error. """

    if isinstance(lst, str):
        lst = [lst]

    content = []  # flattened list of files
    for element in lst:
        if "*" in element:
            globbed = glob.glob(element)
            if len(globbed) > 0:
                content.extend(globbed)
            else:
                raise ValueError(f"No files could be found for pattern: '{element}'")
        else:
            content.append(element)

    return content


def replace_quotes(string):
    """ Replace single quotes with double quotes in a string (such as from the pprint utility to make a valid json file) """

    in_string = False
    for i, letter in enumerate(string):

        if letter == "\"":
            in_string = not in_string  # reverse in_string flag

        elif letter == "'" and in_string is False:  # do not replace single quotes in strings
            string = string[:i] + "\"" + string[i + 1:]  # replace single quote with double quote

    return string


def resize_text(txt_frame, max_size=18):
    """ Resize text to fit the textbox.

    Parameters
    ----------
    txt_frame : pptx.text.text.TextFrame
        The text frame to be resized.
    max_size : int, default 18
        The maximum fontsize of the text.
    """

    # Find all fonts and fit text to the box
    fonts = matplotlib.font_manager.findSystemFonts()

    # Find a font that works
    for font in fonts:
        try:
            txt_frame.fit_text(font_file=font, max_size=max_size)  # some fonts return a 'cannot unpack non-iterable NoneType object'-error
            break
        except Exception:
            pass


###############################################################################
# ------------------- Classes for building presentation --------------------- #
###############################################################################

class PowerPointReport():

    def __init__(self, template=None, size="standard", verbosity=0):
        """ Initialize a presentation object using an existing presentation (template) or from scratch (default) """

        self.template = template

        if template is None:
            self.size = size

        self.setup_logger(verbosity)

        self.logger.info("Initializing presentation")
        self.initialize_presentation()

    def setup_logger(self, verbosity=1):
        """ Setup a logger for the class in self.logger """

        self.logger = logging.getLogger(self.__class__.__name__)

        # Setup formatting of handler
        H = logging.StreamHandler(sys.stdout)
        simple_formatter = logging.Formatter("[%(levelname)s] %(message)s")
        debug_formatter = logging.Formatter("[%(levelname)s] [%(name)s:%(funcName)s] %(message)s")

        # Set verbosity and formatting
        if verbosity == 0:
            self.logger.setLevel(logging.ERROR)
            H.setFormatter(simple_formatter)
        elif verbosity == 1:
            self.logger.setLevel(logging.INFO)
            H.setFormatter(simple_formatter)
        elif verbosity == 2:
            self.logger.setLevel(logging.DEBUG)
            H.setFormatter(debug_formatter)
        else:
            raise ValueError("Verbosity must be 0, 1 or 2.")

        self.logger.addHandler(H)

    def initialize_presentation(self):
        """ Initialize a presentation from scratch. Sets the self._prs and self._slides attributes."""

        self._prs = Presentation(self.template)

        # Set size of the presentation (if not given by a template)
        if self.template is None:
            self.set_size(self.size)  # size is not set if template was given

        # Get ready to add slides
        self._slides = []   # a list of SlidePlus objects

    def set_size(self, size):
        """
        Set the size of the presentation.

        Parameters
        ----------
        size : str or tuple of float
            Size of the presentation. Can be "standard", "widescreen", "a4-portait" or "a4-landscape". Can also be a tuple of numbers indicating (height, width) in cm.
        """

        if isinstance(size, tuple):
            if len(size) != 2:
                raise ValueError("Size tuple must be of length 2.")

            h, w = Cm(size[0]), Cm(size[1])

        elif size == "standard":
            h, w = Cm(19.05), Cm(25.4)

        elif size == "widescreen":
            h, w = Cm(19.05), Cm(33.867)

        elif size == "a4-portrait":
            h, w = Cm(27.517), Cm(19.05)

        elif size == "a4-landscape":
            h, w = Cm(19.05), Cm(27.517)

        else:
            raise ValueError("Invalid size given. Choose from: 'standard', 'widescreen', 'a4-portrait', 'a4-landscape' or a tuple of floats.")

        self._prs.slide_height = h
        self._prs.slide_width = w

    def add_title_slide(self, title, layout=0, subtitle=None):
        """
        Add a title slide to the presentation.

        Parameters
        ----------
        title : str
            Title of the slide.
        layout : int, default 0
            The layout of the slide. The first layout (0) is usually the default title slide.
        subtitle : str, optional
            Subtitle of the slide if the layout has a suptitle placeholder.
        """

        self.add_slide(title=title, slide_layout=layout)
        slide = self._slides[-1]

        # Fill placeholders
        if subtitle is not None:
            if len(slide.placeholders) == 2:
                slide.placeholders[1].text = subtitle

    def add_slide(self,
                  content=[],
                  title=None,
                  slide_layout=1,
                  content_layout="grid",
                  outer_margin=2,
                  inner_margin=1,
                  n_columns=2,
                  width_ratios=None,
                  height_ratios=None,
                  notes=None,
                  split=False,
                  ):
        """
        Add a slide to the presentation.

        Parameters
        ----------
        content : list of str
            List of content to be added to the slide. Can be either a path to a file or a string.
        title : str, optional
            Title of the slide.
        slide_layout : int, default 1
            Layout of the slide.
        content_layout : str, default "grid"
            Layout of the slide. Can be "grid", "vertical" or "horizontal". Can also be an array of integers indicating the layout of the slide.
        outer_margin : float, default 2
            Outer margin of the slide (in cm).
        inner_margin : float, default 1
            Inner margin of the slide elements (in cm).
        n_columns : int, default 2
            Number of columns in the layout in case of "grid" layout.
        width_ratios : list of float, optional
            Width of the columns in case of "grid" layout.
        height_ratios : list of float, optional
            Height of the rows in case of "grid" layout.
        notes : str, optional
            Notes for the slide. Can be either a path to a text file or a string.
        split : bool, default False
            Split the content into multiple slides.
        """

        # Get all parameters
        parameters = locals()

        # Establish content
        if isinstance(content, str):
            content = [content]

        # If split is false, content should be contained in one slide
        if split is False:
            content = [content]
        else:
            content = glob_files(content)

        # Create slide(s)
        for slide_content in content:

            # How many slides are already in the presentation?
            n_slides = len(self._slides)
            self.logger.info("Adding slide {}".format(n_slides + 1))

            # Setup an empty slide
            slide = self.setup_slide(slide_layout)
            slide.add_parameters(parameters)
            slide.content = slide_content

            # Glob files
            slide._content = glob_files(slide_content)  # internal extension of content
            self.logger.debug("Final content: {}".format(slide._content))

            # Set title of slide
            slide.set_title(title)

            # Add content to slide
            if len(slide_content) > 0:
                slide.set_layout_matrix()  # Find the layout of the slide
                slide.create_boxes()       # Create boxes based on layout
                slide.fill_boxes()         # Fill boxes with content

            # Add notes to slide
            slide.add_notes()

    def setup_slide(self, slide_layout):
        """ Initialize an empty slide with a given layout. """

        layout_obj = self._prs.slide_layouts[slide_layout]
        slide_obj = self._prs.slides.add_slide(layout_obj)

        slide = Slide(slide_obj)
        slide.logger = self.logger

        slide._slide_height = self._prs.slide_height
        slide._slide_width = self._prs.slide_width
        self._slides.append(slide)

        return slide

    def add_borders(self):
        """ Add borders of all content boxes. Useful for debugging layouts."""

        for slide in self._slides:
            if hasattr(slide, "_boxes"):
                for box in slide._boxes:
                    box.add_border()

    def remove_borders(self):
        """ Remove borders (is there are any) of all content boxes. """

        for slide in self._slides:
            if hasattr(slide, "_boxes"):
                for box in slide._boxes:
                    box.remove_border()

    def get_config(self, full=False):
        """
        Return a dictionary with the configuration of the presentation

        Parameters
        ----------
        full : bool, default False
            If True, return the full configuration of the presentation. If False, only return the non-default values.

        Returns
        -------
        config : dict
            Dictionary with the configuration of the presentation.
        """

        # Collect upper-level config of presentation
        config = {}
        for key in self.__dict__:
            if key[0] != "_" and key != "logger":  # ignore private attributes and logger
                value = self.__dict__[key]
                if value is not None:  # 'template' can for example be None if no template is used
                    config[key] = value

        # Get default slide config
        defaults = get_default_args(self.add_slide)

        # Get config of each slide
        config["slides"] = []
        for slide in self._slides:
            slide_config = {}
            for key, value in slide.__dict__.items():
                if not key.startswith("_") and key != "logger":  # ignore private attributes and logger
                    value = slide.__dict__[key]

                    if isinstance(value, bool):
                        value_converted = str(value)  # convert bool to str to make it json-compatible
                    else:
                        value_converted = value

                    slide_config[key] = value_converted

                    if full is False:
                        if value == defaults[key]:
                            del slide_config[key]
                        elif isinstance(value, list) and len(value) == 0:  # content can be an empty list
                            del slide_config[key]

            config["slides"].append(slide_config)

        return config

    def write_config(self, filename):
        """ Write the configuration of the presentation to a json-formatted file. """

        config = self.get_config()

        # Get pretty printed config
        pp = pprint.PrettyPrinter(compact=True, sort_dicts=False, width=120)
        config_json = pp.pformat(config)
        config_json = replace_quotes(config_json)
        config_json = re.sub(r"\"\n\s+\"", "", config_json)  # strings are not allowed to split over multiple lines

        with open(filename, "w") as f:
            f.write(config_json)

    def from_config(self, config):
        """
        Fill a presentation from a configuration dictionary.

        Parameters
        ----------
        config : str or dict
            A path to a configuration file or a dictionary containing the configuration (such as from Report.get_config()).
        """

        # Load config from file if necessary
        if isinstance(config, str):
            with open(config, "r") as f:
                try:
                    config = json.load(f)
                except Exception as e:
                    raise ValueError("Could not load config file from {}. The error was: {}".format(config, e))

        # Set upper attributes
        upper_keys = config.keys()
        for key in upper_keys:
            if key != "slides":
                setattr(self, key, config[key])

                if key == "split":
                    self.split = bool(config[key])  # convert input string to bool

        # Initialize presentation
        self.initialize_presentation()

        # Fill in slides with information from slide config
        for slide_dict in config["slides"]:
            self.add_slide(**slide_dict)  # add all options from slide config

    def save(self, filename, show_borders=False, pdf=False):
        """
        Save the presentation to a file.

        Parameters
        ----------
        filename : str
            Filename of the presentation.
        show_borders : bool, default False
            Show borders of the content boxes. Is useful for debugging layouts.
        pdf : bool, default False
            Additionally save the presentation as a pdf file with the same basename as <filename>.
        """

        if show_borders is True:
            self.add_borders()

        self.logger.info("Saving presentation to '" + filename + "'")

        # Warning if filename does nto end with .pptx
        if not filename.endswith(".pptx"):
            self.logger.warning("Filename does not end with '.pptx'. This might cause problems when opening the presentation.")

        self._prs.save(filename)

        # Remove borders again
        if show_borders is True:
            self.remove_borders()  # Remove borders again

        # Save presentation as pdf
        if pdf:

            self.logger.info("Additionally saving presentation as .pdf")

            # Check if libreoffice is installed
            is_installed = False
            try:
                self.logger.debug("Checking if libreoffice is installed...")
                result = subprocess.run(["libreoffice", "--version"], stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
                self.logger.debug("Version of libreoffice: " + result.stdout.rstrip())
                is_installed = True

            except FileNotFoundError:
                self.logger.error("Option 'pdf' is set to True, but LibreOffice could not be found on path. Please install LibreOffice to save presentations as pdf.")

            # Save presentation as pdf
            if is_installed:

                outdir = os.path.dirname(filename)
                outdir = "." if outdir == "" else outdir  # outdir cannot be empty

                cmd = f"libreoffice --headless --invisible --convert-to pdf --outdir {outdir} {filename}"
                self.logger.debug("Running command: " + cmd)

                process = subprocess.Popen(cmd, shell=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
                while process.poll() is None:
                    line = process.stdout.readline().rstrip()
                    if line != "":
                        self.logger.debug("Command output: " + line)


# ------------------------------------------------------------------------------

class Slide():
    """ An internal class for creating slides. """

    def __init__(self, slide):

        self._slide = slide  # Slide object from python-pptx
        self._boxes = []   # Boxes in the slide
        self.logger = None

    def set_title(self, title):
        """ Set the title of the slide. """

        if title is not None:

            if self._slide.shapes.title is None:
                self.logger.warning("Could not set title of slide. The slide does not have a title box.")
            else:
                self._slide.shapes.title.text = title

    def add_parameters(self, parameters):
        """ Add parameters to the slide. """

        for key in parameters:
            if key != "self":
                setattr(self, key, parameters[key])

    def add_notes(self):
        """ Add notes to the slide. """

        if self.notes is not None:
            if os.path.exists(self.notes):
                with open(self.notes, "r") as f:
                    notes_string = f.read()
            else:
                notes_string = self.notes

            self._slide.notes_slide.notes_text_frame.text = notes_string

    def set_layout_matrix(self):
        """ Get the content layout matrix for the slide. """

        # Get variables from self
        layout = self.content_layout
        n_elements = len(self._content)
        n_columns = self.n_columns

        # Get layout matrix depending on "layout" variable
        if layout == "grid":
            n_columns = min(n_columns, n_elements)  # number of columns cannot be larger than number of elements
            n_rows = int(np.ceil(n_elements / n_columns))  # number of rows to fit elements
            n_total = n_rows * n_columns

            intarray = list(range(n_elements))
            intarray.extend([np.nan] * (n_total - n_elements))

            layout_matrix = np.array(intarray).reshape((n_rows, n_columns))

        elif layout == "vertical":
            layout_matrix = np.array(list(range(n_elements))).reshape((n_elements, 1))

        elif layout == "horizontal":
            layout_matrix = np.array(list(range(n_elements))).reshape((1, n_elements))

        else:
            layout_matrix = self._validate_layout(layout)  # check if layout is a valid matrix

        self._layout_matrix = layout_matrix

    @staticmethod
    def _validate_layout(layout_matrix):
        """ Validate the given layout matrix. """
        # TODO: check if layout is a valid matrix

        layout_matrix = np.array(layout_matrix)

        return layout_matrix

    def create_boxes(self):
        """ Create boxes for the slide dependent on the intrnal layout matrix. """

        layout_matrix = self._layout_matrix
        nrows, ncols = layout_matrix.shape

        # Convert margins from cm to pptx units
        outer_margin_unit = Cm(self.outer_margin)
        inner_margin_unit = Cm(self.inner_margin)
        left_margin = right_margin = outer_margin_unit
        top_margin = bottom_margin = outer_margin_unit

        # Add to top margin based on title
        if self._slide.shapes.title.text != "":
            top_margin = self._slide.shapes.title.top + self._slide.shapes.title.height + outer_margin_unit
        else:
            top_margin = outer_margin_unit
            sp = self._slide.shapes.title.element
            sp.getparent().remove(sp)

        # How many columns and rows are there?
        n_rows, n_cols = layout_matrix.shape

        # Get total height and width of pictures
        total_width = self._slide_width - left_margin - right_margin - (n_cols - 1) * inner_margin_unit
        total_height = self._slide_height - top_margin - bottom_margin - (n_rows - 1) * inner_margin_unit

        # Get column widths and row heights
        if self.width_ratios is None:
            widths = (np.ones(ncols) / ncols) * total_width
        else:
            widths = np.array(self.width_ratios) / sum(self.width_ratios) * total_width

        if self.height_ratios is None:
            heights = (np.ones(nrows) / nrows) * total_height
        else:
            heights = np.array(self.height_ratios) / sum(self.height_ratios) * total_height

        # Box coordinates
        box_numbers = layout_matrix[~np.isnan(layout_matrix)].flatten()
        box_numbers = sorted(set(box_numbers))  # unique box numbers
        box_numbers = [box_number for box_number in box_numbers if box_number >= 0]  # remove negative numbers
        for i in box_numbers:

            # Get column and row number
            coordinates = np.argwhere(layout_matrix == i)

            # Get upper left corner of box
            row, col = coordinates[0]
            left = left_margin + np.sum(widths[:col]) + col * inner_margin_unit
            top = top_margin + np.sum(heights[:row]) + row * inner_margin_unit

            # Get total width and height of box (can span multiple columns and rows)
            width = 0
            height = 0
            rows = set(coordinates[:, 0])
            for row in rows:
                height += heights[row]
            height += (len(rows) - 1) * inner_margin_unit  # add inner margins between rows

            cols = set(coordinates[:, 1])
            for col in cols:
                width += widths[col]
            width += (len(cols) - 1) * inner_margin_unit  # add inner margins between columns

            #  Create box
            self.add_box((left, top, width, height))

    def add_box(self, coordinates):
        """
        Add a box to the slide.

        Parameters
        ----------
        coordinates : tuple
            Coordinates containing (left, top, width, height) of the box (in pptx units).
        """

        box = Box(self._slide, coordinates)
        box.logger = self.logger  # share logger with box
        self._boxes.append(box)

    def fill_boxes(self):
        """ Fill the boxes with the elements in self.content """

        for i, element in enumerate(self._content):
            self._boxes[i].fill(element)


class Box():
    """ A box is a constrained area of the slide which contains a single element e.g. text, a picture, a table, etc. """

    def __init__(self, slide, coordinates):
        """
        Initialize a box.

        Parameters
        ----------
        slide : pptx slide object
            The slide on which the box is located.
        coordinates : tuple
            Coordinates containing (left, top, width, height) of the box (in pptx units).
        """

        self.slide = slide
        self.logger = None

        # Bounds of the box
        self.left = int(coordinates[0])
        self.top = int(coordinates[1])
        self.width = int(coordinates[2])
        self.height = int(coordinates[3])

        # Initialize bounds of the content (can be smaller than the box)
        self.content = None
        self.content_left = self.left
        self.content_top = self.top
        self.content_width = self.width
        self.content_height = self.height

        self.border = None  # border object of the box

    def add_border(self):
        """ Adds a border shape of box to make debugging easier """

        if self.border is None:
            self.border = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.left, self.top, self.width, self.height)
            self.border.fill.background()
            self.border.line.color.rgb = RGBColor(0, 0, 0)  # black
            self.border.line.width = Pt(1)

    def remove_border(self):
        """ Removes the border shape """

        if self.border is not None:
            self.border._sp.getparent().remove(self.border._sp)
            self.border = None  # reset border

    def fill(self, content):
        """
        Fill the box with content. The function estimates type of content is given.

        Parameters
        ----------
        content : str
            The element to be added to the box.
        """

        self.content = content

        # Find out what type of content it is
        content_type = self._get_content_type(content)

        if content_type == "pdf":
            filename = self.convert_pdf(content)
            self.fill_image(filename)
            os.remove(filename)

        elif content_type == "image":
            self.fill_image(content)

        elif content_type == "textfile":
            with open(content) as f:
                text = f.read()
            self.fill_text(text)

        elif content_type == "text":
            self.fill_text(content)

        else:
            pass

    @staticmethod
    def _get_content_type(content):
        """ Determine the type of content. """

        if isinstance(content, str):
            if os.path.isfile(content):

                # Find out the content of the file
                try:
                    with open(content) as f:
                        _ = f.read()
                    return "textfile"

                except UnicodeDecodeError:

                    if content.endswith(".pdf"):
                        return "pdf"
                    else:
                        return "image"
            else:
                return "text"
        else:
            t = type(content)
            raise ValueError(f"Content of type '{t}' cannot be added to slide")

    def convert_pdf(self, pdf):
        """ Convert a pdf file to a png file. """

        # Create temporary file
        temp_name = next(tempfile._get_candidate_names()) + ".png"
        temp_dir = tempfile.gettempdir()
        temp_file = os.path.join(temp_dir, temp_name)
        self.logger.debug(f"Converting pdf to temporary png at: {temp_file}")

        # Convert pdf to png
        doc = fitz.open(pdf)
        page = doc.load_page(0)
        pix = page.get_pixmap()
        pix.save(temp_file)

        return temp_file

    def fill_image(self, filename):
        """ Fill the box with an image. """

        # Find out the size of the image
        self._adjust_image_size(filename)
        self._adjust_image_position()  # adjust image position to middle of box

        # Add image
        self.logger.debug("Adding image to slide from file: " + filename)
        self.picture = self.slide.shapes.add_picture(filename, self.content_left, self.content_top, self.content_width, self.content_height)

    def _adjust_image_size(self, filename):
        """
        Adjust the size of the image to fit the box.

        Parameters
        ----------
        filename : str
            Path to the image file.
        """

        # Find out the size of the image
        im = Image.open(filename)
        im_width, im_height = im.size

        box_width = self.width
        box_height = self.height

        im_ratio = im_width / im_height  # >1 for landscape, <1 for portrait
        box_ratio = box_width / box_height

        # width is the limiting factor; height will be smaller than box_height
        if box_ratio < im_ratio:  # box is wider than image; height will be limiting
            self.content_width = box_width
            self.content_height = box_width * im_height / im_width  # maintain aspect ratio

        # height is the limiting factor; width will be smaller than box_width
        else:
            self.content_width = box_height * im_width / im_height  # maintain aspect ratio
            self.content_height = box_height

    def _adjust_image_position(self):
        """ Adjust the position of the image to be in the middle of the box. """

        self.content_left = self.left + (self.width - self.content_width) / 2
        self.content_top = self.top + (self.height - self.content_height) / 2

    def fill_text(self, text):
        """
        Fill the box with text.

        Parameters
        ----------
        text : str
            The text to be added to the box.
        """

        txt_box = self.slide.shapes.add_textbox(self.left, self.top, self.width, self.height)
        txt_frame = txt_box.text_frame

        txt_frame.add_paragraph()
        p = txt_frame.paragraphs[0]
        p.text = text
        txt_frame.word_wrap = True
        txt_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Try to fit text to the box
        resize_text(txt_frame)
        txt_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # An additional step of resizing text to fit the box
