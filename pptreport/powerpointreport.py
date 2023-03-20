import os
import glob
import inspect
import pprint
import json
import re
import subprocess
import logging
import sys
from natsort import natsorted

# Pptx modules
from pptx import Presentation
from pptx.util import Cm

from pptreport.slide import Slide

###############################################################################
# ---------------------------- Helper functions ----------------------------- #
###############################################################################


def flatten_list(lst):
    """ Flatten a list containing both lists and non-lists """

    flat = []
    for element in lst:
        if isinstance(element, list):
            flat.extend(element)
        else:
            flat.append(element)

    return flat


def get_default_args(func):
    signature = inspect.signature(func)
    defaults = {k: v.default for k, v in signature.parameters.items() if v.default is not inspect.Parameter.empty}

    return defaults


def glob_files(lst):
    """ Glob files in a list of strings which might contain "*". If no files are found for glob, raise an error. """

    if isinstance(lst, str):
        lst = [lst]

    content = []  # flattened list of files
    files = []  # names of files in the list
    for element in lst:
        if "*" in element:
            globbed = glob.glob(element)
            if len(globbed) > 0:
                for file in globbed:
                    files.append(file)
                    content.append(file)
            else:
                raise ValueError(f"No files could be found for pattern: '{element}'")
        else:
            content.append(element)

    # Get the locations of files in the list
    file_locations = []
    for i, string in enumerate(content):
        if string in files:
            file_locations.append(i)

    # Sort files using natural sorting
    file_list = [content[i] for i in file_locations]
    file_list = natsorted(file_list)

    # Return files to content in the correct order
    for idx, string in zip(file_locations, file_list):
        content[idx] = string

    return content


def get_files_in_dir(directory):
    """ Get all files in the given directory including the path prefix """

    files = os.listdir(directory)
    files = [os.path.join(directory, file) for file in files]

    return files


def fill_dict(d1, d2):
    """ Fill the keys of d1 with the values of d2 if they are not already present in d1.

    Returns
    --------
    None
        d1 is updated in place.
    """

    for key, value in d2.items():
        if key not in d1:
            d1[key] = value


def replace_quotes(string):
    """ Replace single quotes with double quotes in a string (such as from the pprint utility to make a valid json file) """

    in_string = False
    for i, letter in enumerate(string):

        if letter == "\"":
            in_string = not in_string  # reverse in_string flag

        elif letter == "'" and in_string is False:  # do not replace single quotes in strings
            string = string[:i] + "\"" + string[i + 1:]  # replace single quote with double quote

    return string


###############################################################################
# -------------------- Class for building presentation ---------------------- #
###############################################################################

class PowerPointReport():
    """ Class for building a PowerPoint presentation """

    _default_slide_parameters = {
        "title": None,
        "slide_layout": 1,
        "content_layout": "grid",
        "content_alignment": "center",
        "outer_margin": 2,
        "inner_margin": 1,
        "left_margin": None,
        "right_margin": None,
        "top_margin": None,
        "bottom_margin": None,
        "n_columns": 2,
        "width_ratios": None,
        "height_ratios": None,
        "notes": None,
        "split": False,
        "show_filename": False,
        "filename_alignment": "center",
        "filename_path": False,
        "fill_by": "row",
        "remove_placeholders": False,
        "fontsize": None
    }

    def __init__(self, template=None, size="standard", verbosity=0):
        """ Initialize a presentation object using an existing presentation (template) or from scratch (default) """

        self.template = template
        if template is None:
            self.size = size

        self.global_parameters = None
        self.setup_logger(verbosity)

        self.logger.info("Initializing presentation")
        self._initialize_presentation()

    def setup_logger(self, verbosity=1):
        """
        Setup a logger for the class.

        Parameters
        ----------
        verbosity : int, default 1
            The verbosity of the logger. 0: ERROR, 1: INFO, 2: DEBUG

        Returns
        -------
        None
            self.logger is set.
        """

        self.logger = logging.getLogger(self.__class__.__name__)

        # Setup formatting of handler
        H = logging.StreamHandler(sys.stdout)
        simple_formatter = logging.Formatter("[%(levelname)s] %(message)s")
        debug_formatter = logging.Formatter("[%(levelname)s] [%(name)s:%(funcName)s] %(message)s")

        # Set verbosity and formatting
        if verbosity == 0:
            self.logger.setLevel(logging.ERROR)
            H.setFormatter(simple_formatter)
        elif verbosity == 1:
            self.logger.setLevel(logging.INFO)
            H.setFormatter(simple_formatter)
        elif verbosity == 2:
            self.logger.setLevel(logging.DEBUG)
            H.setFormatter(debug_formatter)
        else:
            raise ValueError("Verbosity must be 0, 1 or 2.")

        self.logger.addHandler(H)

    def _initialize_presentation(self):
        """ Initialize a presentation from scratch. Sets the self._prs and self._slides attributes."""

        self._prs = Presentation(self.template)

        # Get ready to collect configuration
        self._config_dict = {}  # configuration dictionary

        # Set size of the presentation (if not given by a template)
        if self.template is None:
            self.set_size(self.size)  # size is not set if template was given

        # Get ready to add slides
        self._slides = []   # a list of Slide objects

        # Add info to config dict
        if self.template is not None:
            self._config_dict["template"] = self.template
        else:
            self._config_dict["size"] = self.size

    def add_global_parameters(self, parameters):
        """ Add global parameters to the presentation """

        # Test that parameters is a dict
        if not isinstance(parameters, dict):
            raise TypeError("Parameters must be a dict.")

        # Save parameters to self
        self.global_parameters = parameters  # for writing to config file

        # Overwrite default parameters
        for k, v in parameters.items():
            if k not in self._default_slide_parameters:
                raise ValueError(f"Parameter '{k}' is not a valid parameter for slide.")
            else:
                self._default_slide_parameters[k] = v

        # Add to internal config dict
        self._config_dict["global_parameters"] = parameters

    def add_to_config(self, parameters):
        """ Add the slide parameters to the config file """

        if "slides" not in self._config_dict:
            self._config_dict["slides"] = []

        self._config_dict["slides"].append(parameters)

    def set_size(self, size):
        """
        Set the size of the presentation.

        Parameters
        ----------
        size : str or tuple of float
            Size of the presentation. Can be "standard", "widescreen", "a4-portait" or "a4-landscape". Can also be a tuple of numbers indicating (height, width) in cm.
        """

        if isinstance(size, tuple):
            if len(size) != 2:
                raise ValueError("Size tuple must be of length 2.")

            h, w = Cm(size[0]), Cm(size[1])

        elif size == "standard":
            h, w = Cm(19.05), Cm(25.4)

        elif size == "widescreen":
            h, w = Cm(19.05), Cm(33.867)

        elif size == "a4-portrait":
            h, w = Cm(27.517), Cm(19.05)

        elif size == "a4-landscape":
            h, w = Cm(19.05), Cm(27.517)

        else:
            raise ValueError("Invalid size given. Choose from: 'standard', 'widescreen', 'a4-portrait', 'a4-landscape' or a tuple of floats.")

        self._prs.slide_height = h
        self._prs.slide_width = w

    def format_parameters(self, parameters):
        """ Checks and formats specific slide parameters to the correct type. """

        parameters = parameters.copy()  # Make a copy to not change the original dict

        # Format "n_columns" to int
        if "n_columns" in parameters:
            try:
                parameters["n_columns"] = int(parameters["n_columns"])
            except ValueError:
                raise ValueError(f"Could not convert 'n_columns' parameter to int. The given value is: '{parameters['n_columns']}'. Please use an integer.")

        # Format boolean parameters to bool
        bool_parameters = [key for key, value in self._default_slide_parameters.items() if isinstance(value, bool)]
        for param in bool_parameters:
            if param in parameters:
                value = parameters[param]
                if isinstance(value, str):
                    if value.lower() in ["true", "1", "t", "y", "yes"]:
                        parameters[param] = True
                    elif value.lower() in ["false", "0", "f", "n", "no"]:
                        parameters[param] = False
                    else:
                        raise ValueError(f"Could not convert '{param}' parameter to bool. The given value is: '{value}'. Please use 'True' or 'False'.")

        return parameters

    # ------------------------------------------------------ #
    # ------------- Functions for adding slides ------------ #
    # ------------------------------------------------------ #

    def add_title_slide(self, title, layout=0, subtitle=None):
        """
        Add a title slide to the presentation.

        Parameters
        ----------
        title : str
            Title of the slide.
        layout : int, default 0
            The layout of the slide. The first layout (0) is usually the default title slide.
        subtitle : str, optional
            Subtitle of the slide if the layout has a suptitle placeholder.
        """

        self.add_slide(title=title, slide_layout=layout)
        slide = self._slides[-1]

        # Fill placeholders
        if subtitle is not None:
            if len(slide.placeholders) == 2:
                slide.placeholders[1].text = subtitle

    def add_slide(self,
                  content=None,
                  grouped_content=None,
                  title=None,
                  slide_layout=None,
                  content_layout=None,
                  content_alignment=None,
                  outer_margin=None,
                  inner_margin=None,
                  left_margin=None,
                  right_margin=None,
                  top_margin=None,
                  bottom_margin=None,
                  n_columns=None,
                  width_ratios=None,
                  height_ratios=None,
                  notes=None,
                  split=None,
                  show_filename=None,
                  filename_alignment=None,
                  filename_path=None,
                  fill_by=None,
                  remove_placeholders=None,
                  fontsize=None,
                  ):
        """
        Add a slide to the presentation.

        Parameters
        ----------
        content : list of str
            List of content to be added to the slide. Can be either a path to a file or a string.
        grouped_content : list of str
            List of grouped content to be added to the slide. The groups are identified by the regex groups of each element in the list.
        title : str, optional
            Title of the slide.
        slide_layout : int or str, default 1
            Layout of the slide. If an integer, it is the index of the layout. If a string, it is the name of the layout.
        content_layout : str, default "grid"
            Layout of the slide. Can be "grid", "vertical" or "horizontal". Can also be an array of integers indicating the layout of the slide.
        content_alignment : str, default "center"
            Alignment of the content. Can be combinations of "upper", "lower", "left", "right" and "center". Examples: "upper left", "center", "lower right".
            The default is "center", which will align the content centered both vertically and horizontally.
        outer_margin : float, default 2
            Outer margin of the slide (in cm).
        inner_margin : float, default 1
            Inner margin of the slide elements (in cm).
        left_margin / right_margin : float, optional
            Left and right margin of the slide elements (in cm). Can be used to overwrite outer_margin for left/right/both dependent on which are given.
        top_margin / bottom_margin : float, optional
            Top and bottom margin of the slide elements (in cm). Can be used to overwrite outer_margin for top/bottom/both dependent on which are given.
        n_columns : int, default 2
            Number of columns in the layout in case of "grid" layout.
        width_ratios : list of float, optional
            Width of the columns in case of "grid" layout.
        height_ratios : list of float, optional
            Height of the rows in case of "grid" layout.
        notes : str, optional
            Notes for the slide. Can be either a path to a text file or a string.
        split : bool or int, default False
            Split the content into multiple slides. If True, the content will be split into one-slide-per-element. If an integer, the content will be split into slides with that many elements per slide.
        show_filename : bool, default False
            Filenames for images. If True, the filename of the image will be displayed above the image.
        filename_alignment : str, default "center"
            Horizontal alignment of the filename. Can be "left", "right" and "center".
            The default is "center", which will align the content centered horizontally.
        filename_path : bool, default False
            Whether to show the full path of the filename or just the filename. Default is False, which will only show the filename.
        fill_by : str, default "row"
            If slide_layout is grid or custom, choose to fill the grid row-by-row or column-by-column. 'fill_by' can be "row" or "column".
        remove_placeholders : str, default False
            Whether to remove empty placeholders from the slide, e.g. if title is not given. Default is False; to keep all placeholders. If True, empty placeholders will be removed.
        fontsize : float, default None
            Fontsize of text content. If None, the fontsize is automatically determined to fit the text in the textbox.
        """

        # Get input parameters; all function defaults are None to distinguish between given arguments and global defaults
        parameters = locals()
        parameters = {k: v for k, v in parameters.items() if v is not None}
        parameters.pop("self")
        parameters = self._check_add_slide_input(parameters)
        self.logger.debug(f"Input parameters: {parameters}")

        # If input was None, replace with default parameters from upper presentation
        fill_dict(parameters, self._default_slide_parameters)
        self.add_to_config(parameters)
        self.logger.debug("Final slide parameters: {}".format(parameters))

        # Check validity and format parameters before creating slides
        parameters = self.format_parameters(parameters)

        # Add slides dependent on content type
        if "grouped_content" in parameters:

            content_per_group = self._get_paired_content(parameters["grouped_content"])

            # Create one slide per group
            for group, content in content_per_group.items():
                slide = self._setup_slide(parameters)
                slide.title = f"Group: {group}" if slide.title is None else slide.title
                slide.content = content
                slide._fill_slide()

        else:

            content = self._get_content(parameters)

            # Create slide(s)
            for slide_content in content:

                # Setup an empty slide
                slide = self._setup_slide(parameters)
                slide.content = slide_content
                slide._fill_slide()  # Fill slide with content

    def _check_add_slide_input(self, parameters):
        """ Check the format of the input parameters for the slide and return an updated dictionary. """

        # Establish if content or grouped_content was given
        if "content" in parameters and "grouped_content" in parameters:
            raise ValueError("Invalid input. Both 'content' and 'grouped_content' were given - please give only one input type.")

        # Check format of content
        # if "content" in parameters:

        # If split is given, content should be given
        if parameters.get("split", False) is not False and len(parameters.get("content", [])) == 0:
            raise ValueError("Invalid input. 'split' is given, but 'content' is empty")

        return parameters

    def _setup_slide(self, parameters):
        """ Initialize an empty slide with a given layout. """

        # How many slides are already in the presentation?
        n_slides = len(self._slides)
        self.logger.info("Adding slide {}".format(n_slides + 1))

        # Add slide to python-pptx presentation
        slide_layout = parameters.get("slide_layout", 0)
        layout_obj = self._get_slide_layout(slide_layout)
        slide_obj = self._prs.slides.add_slide(layout_obj)

        # Add slide to list of slides in internal object
        slide = Slide(slide_obj, parameters)
        slide.logger = self.logger

        # Add information from presentation to slide
        slide._default_parameters = self._default_slide_parameters
        slide._slide_height = self._prs.slide_height
        slide._slide_width = self._prs.slide_width

        self._slides.append(slide)

        return slide

    def _get_slide_layout(self, slide_layout):
        """ Get the slide layout object from a given layout. """

        if isinstance(slide_layout, int):
            try:
                layout_obj = self._prs.slide_layouts[slide_layout]
            except IndexError:
                n_layouts = len(self._prs.slide_layouts)
                raise IndexError(f"Layout index {slide_layout} not found in slide master. The number of slide layouts is {n_layouts} (the maximum index is {n_layouts-1})")

        elif isinstance(slide_layout, str):
            try:
                layout_obj = self._prs.slide_layouts.get_by_name(slide_layout)
            except KeyError:
                raise KeyError(f"Layout named '{slide_layout}' not found in slide master.")
        else:
            raise TypeError("Layout should be an integer or a string.")

        return layout_obj

    def _get_content(self, parameters):
        """ Get slide content based on input parameters. """

        # Establish content
        content = parameters.get("content", [])
        if isinstance(content, str):
            content = [content]

        # Expand content files
        content = glob_files(content)

        # If split is false, content should be contained in one slide
        if parameters["split"] is False:
            content = [content]
        else:
            if len(content) == 0:
                raise ValueError("Split is True, but 'content' is empty.")
            else:
                content = glob_files(content)
                if isinstance(parameters["split"], int):
                    content = [content[i:i + parameters["split"]] for i in range(0, len(content), parameters["split"])]

        return content

    def _get_paired_content(self, raw_content):
        """ Get content per group from a list of regex patterns.

        Parameters
        ----------
        raw_content : list of str
            List of regex patterns. Each pattern should contain one group.

        Returns
        -------
        content_per_group : dict
            Dictionary with group names as keys and lists of content as values.
        """

        # Search for regex groups
        group_content = {}  # dict of lists of content input
        for i, pattern in enumerate(raw_content):
            group_content[i] = {}

            self.logger.debug(f"Finding files for pattern: {pattern}")

            # Establish folder and all files in it
            dirname = os.path.dirname(pattern)
            dirname = "." if dirname == "" else dirname
            files = get_files_in_dir(dirname)

            # Find all files that match the regex
            try:
                pattern_compiled = re.compile(pattern)
            except re.error:
                raise ValueError(f"Invalid regex: {pattern}")

            # Find all files that match the regex
            for fil in files:

                m = pattern_compiled.match(fil)
                if m:  # if there was a match

                    groups = m.groups()
                    if len(groups) == 0:
                        raise ValueError(f"Regex {pattern} does not contain any groups.")
                    elif len(groups) > 1:
                        raise ValueError(f"Regex {pattern} contains more than one group.")
                    group = groups[0]

                    # Save the file to the group
                    group_content[i][group] = fil

        # Collect all groups found
        all_regex_groups = sum([list(d.keys()) for d in group_content.values()], [])  # flatten list of lists
        all_regex_groups = natsorted(set(all_regex_groups))
        self.logger.debug(f"Found groups: {all_regex_groups}")

        # If no groups were found for an element, add strings for each group
        for i in group_content:
            if len(group_content[i]) == 0:
                for group in all_regex_groups:
                    group_content[i][group] = raw_content[i]

        # Convert from group per element to element per group
        content_per_group = {group: [group_content[i].get(group, None) for i in group_content] for group in all_regex_groups}

        return content_per_group

    # ------------------------------------------------------------------------ #
    # --------------------- Additional elements on slides -------------------- #
    # ------------------------------------------------------------------------ #

    def add_borders(self):
        """ Add borders of all content boxes. Useful for debugging layouts."""

        for slide in self._slides:
            if hasattr(slide, "_boxes"):
                for box in slide._boxes:
                    box.add_border()

    def remove_borders(self):
        """ Remove borders (is there are any) of all content boxes. """

        for slide in self._slides:
            if hasattr(slide, "_boxes"):
                for box in slide._boxes:
                    box.remove_border()

    # ------------------------------------------------------------------------ #
    # --------------------- Saving / loading presentations ---------------------
    # ------------------------------------------------------------------------ #

    def get_config(self, full=False, expand=False):
        """
        Collect a dictionary with the configuration of the presentation

        Parameters
        ----------
        full : bool, default False
            If True, return the full configuration of the presentation. If False, only return the non-default values.
        expand : bool, default False
            If True, expand the content of each slide to a list of files. If False, keep the content as input including "*" and regex.

        Returns
        -------
        config : dict
            Dictionary with the configuration of the presentation.
        """

        # Get configuration of presentation
        if expand is True:  # Read parameters directly from report object

            config = dict(self.__dict__)

            # Add configuration of each slide
            config["slides"] = []
            for slide in self._slides:
                config["slides"].append(slide.get_config())

            # Remove internal variables
            for key in list(config.keys()):  # list to prevent RuntimeError: dictionary changed size during iteration
                if key.startswith("_") or key == "logger":
                    del config[key]

        else:  # Read parameters from internal config_dict
            config = self._config_dict.copy()

        # Get default slide parameters
        defaults = self._default_slide_parameters

        # Resolve configuration of each slide
        for slide_config in config.get("slides", []):
            for key in list(slide_config.keys()):  # list to prevent RuntimeError: dictionary changed size during iteration
                value = slide_config[key]

                # convert bool to str to make it json-compatible
                if isinstance(value, bool):
                    value_converted = str(value)  # convert bool to str to make it json-compatible
                else:
                    value_converted = value
                slide_config[key] = value_converted

                # Remove default values if full is False
                if full is False:
                    if value == defaults.get(key, None):  # compares to the unconverted value
                        del slide_config[key]
                    elif isinstance(value, list) and len(value) == 0:  # content can be an empty list
                        del slide_config[key]

        return config

    def write_config(self, filename, full=False, expand=False):
        """
        Write the configuration of the presentation to a json-formatted file.

        Parameters
        ----------
        filename : str
            Path to the file to write the configuration to.
        full : bool, default False
            If True, write the full configuration of the presentation. If False, only write the non-default values.
        expand : bool, default False
            If True, expand the content of each slide to a list of files. If False, keep the content as input including "*" and regex.
        """

        config = self.get_config(full=full)

        # Get pretty printed config
        pp = pprint.PrettyPrinter(compact=True, sort_dicts=False, width=120)
        config_json = pp.pformat(config)
        config_json = replace_quotes(config_json)
        config_json = re.sub(r"\"\n\s+\"", "", config_json)  # strings are not allowed to split over multiple lines
        config_json = re.sub(r": None", ": null", config_json)  # Convert to null as None is not allowed in json
        config_json += "\n"  # end with newline

        with open(filename, "w") as f:
            f.write(config_json)

    def from_config(self, config):
        """
        Fill a presentation from a configuration dictionary.

        Parameters
        ----------
        config : str or dict
            A path to a configuration file or a dictionary containing the configuration (such as from Report.get_config()).
        """

        # Load config from file if necessary
        if isinstance(config, str):
            with open(config, "r") as f:
                try:
                    config = json.load(f)
                except Exception as e:
                    raise ValueError("Could not load config file from {}. The error was: {}".format(config, e))

        # Set upper attributes
        upper_keys = config.keys()
        for key in upper_keys:
            if key != "slides":
                setattr(self, key, config[key])

                if key == "split":
                    self.split = bool(config[key])  # convert input string to bool

        # Initialize presentation
        self._initialize_presentation()

        # Set global slide parameters
        if "global_parameters" in config:
            self.add_global_parameters(config["global_parameters"])

        # Fill in slides with information from slide config
        for slide_dict in config["slides"]:
            self.add_slide(**slide_dict)  # add all options from slide config

    def save(self, filename, show_borders=False, pdf=False):
        """
        Save the presentation to a file.

        Parameters
        ----------
        filename : str
            Filename of the presentation.
        show_borders : bool, default False
            Show borders of the content boxes. Is useful for debugging layouts.
        pdf : bool, default False
            Additionally save the presentation as a pdf file with the same basename as <filename>.
        """

        if show_borders is True:
            self.add_borders()

        self.logger.info("Saving presentation to '" + filename + "'")

        # Warning if filename does nto end with .pptx
        if not filename.endswith(".pptx"):
            self.logger.warning("Filename does not end with '.pptx'. This might cause problems when opening the presentation.")

        self._prs.save(filename)

        # Remove borders again
        if show_borders is True:
            self.remove_borders()  # Remove borders again

        # Save presentation as pdf
        if pdf:

            self.logger.info("Additionally saving presentation as .pdf")

            # Check if libreoffice is installed
            is_installed = False
            try:
                self.logger.debug("Checking if libreoffice is installed...")
                result = subprocess.run(["libreoffice", "--version"], stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
                self.logger.debug("Version of libreoffice: " + result.stdout.rstrip())
                is_installed = True

            except FileNotFoundError:
                self.logger.error("Option 'pdf' is set to True, but LibreOffice could not be found on path. Please install LibreOffice to save presentations as pdf.")

            # Save presentation as pdf
            if is_installed:

                outdir = os.path.dirname(filename)
                outdir = "." if outdir == "" else outdir  # outdir cannot be empty

                cmd = f"libreoffice --headless --invisible --convert-to pdf --outdir {outdir} {filename}"
                self.logger.debug("Running command: " + cmd)

                process = subprocess.Popen(cmd, shell=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
                while process.poll() is None:
                    line = process.stdout.readline().rstrip()
                    if line != "":
                        self.logger.debug("Command output: " + line)
