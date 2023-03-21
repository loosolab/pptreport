import numpy as np
import os
from pptx.util import Cm
from pptreport.box import Box
import warnings
from numpy import VisibleDeprecationWarning


class Slide():
    """ An internal class for creating slides. """

    def __init__(self, slide, parameters={}):

        self._slide = slide  # Slide object from python-pptx
        self._boxes = []     # Boxes in the slide
        self.logger = None

        self.add_parameters(parameters)

    def add_parameters(self, parameters):
        """ Add parameters to the slide as internal variables. """

        for key in parameters:
            if key != "self":
                setattr(self, key, parameters[key])

    def get_config(self):
        """ Get the config dictionary for this slide. """

        config = self.__dict__.copy()  # Make a copy to not change the original dict
        for key in list(config):
            if key.startswith("_") or key == "logger":
                del config[key]

        return config

    def set_layout_matrix(self):
        """ Get the content layout matrix for the slide. """

        # Check validity of n_columns
        try:
            self.n_columns = int(self.n_columns)
        except ValueError:
            raise ValueError(f"Could not convert 'n_columns' parameter to int. The given value is: '{self.n_columns}'. Please use an integer.")

        # Get variables from self
        layout = self.content_layout
        n_elements = len(self.content)
        n_columns = self.n_columns

        # Get layout matrix depending on "layout" variable
        if isinstance(layout, str):
            if layout == "grid":

                n_columns = min(n_columns, n_elements)  # number of columns cannot be larger than number of elements
                n_rows = int(np.ceil(n_elements / n_columns))  # number of rows to fit elements
                n_total = n_rows * n_columns

                intarray = list(range(n_elements))
                intarray.extend([np.nan] * (n_total - n_elements))

                if self.fill_by == "row":
                    layout_matrix = np.array(intarray).reshape((n_rows, n_columns))
                elif self.fill_by == "column":
                    layout_matrix = np.array(intarray).reshape((n_columns, n_rows))
                    layout_matrix = layout_matrix.T
                else:
                    raise ValueError(f"Invalid value for 'fill_by' parameter: '{self.fill_by}'. Please use 'row' or 'column'.")

            elif layout == "vertical":
                layout_matrix = np.array(list(range(n_elements))).reshape((n_elements, 1))

            elif layout == "horizontal":
                layout_matrix = np.array(list(range(n_elements))).reshape((1, n_elements))
            else:
                raise ValueError(f"Unknown layout string: '{layout}'. Please use 'grid', 'vertical' or 'horizontal', or a custom matrix.")

        else:  # layout is expected to be a matrix
            layout_matrix = self._validate_layout(layout)  # check if layout is a valid matrix

        self._layout_matrix = layout_matrix

    # ------------------------ Validate options ------------------------#
    def _validate_margins(self):
        """ Check whether the given margins are valid """

        margins = {"outer_margin": self.outer_margin, "inner_margin": self.inner_margin, "left_margin": self.left_margin, "right_margin": self.right_margin,
                   "top_margin": self.top_margin, "bottom_margin": self.bottom_margin}

        for margin, value in margins.items():
            if value is not None:

                # Check whether value is a float
                try:
                    value = float(value)
                    setattr(self, margin, value)
                except ValueError:
                    raise ValueError(f"Could not convert '{margin}' to a float. The given value is: {value}")

                # Check whether value is positive
                if value < 0:
                    raise ValueError(f"Margin '{margin}' cannot be negative. The given value is: {value}")

                # Check upper margin sizes

    def _validate_ratios(self):
        """ Validate the values of width and height ratios """

        parameters = ["width_ratios", "height_ratios"]

        for param in parameters:

            value = getattr(self, param)

            if value is None:
                continue

            # Convert from string to list
            if isinstance(value, str):
                try:
                    value = [v for v in value.split(",")]
                except ValueError:
                    raise ValueError(f"Could not convert '{param}' parameter to list of values. The given value is: '{value}'. Please use a list of values.")

            # Convert from list of strings to list of floats
            try:
                value = [float(v) for v in value]
            except Exception:
                raise ValueError(f"Could not convert '{param}' parameter to list of values. The given value is: '{value}'. Please use a list of values.")

            setattr(self, param, value)  # Set the new value

    @staticmethod
    def _validate_layout(layout_matrix):
        """ Validate the given layout matrix. """
        # TODO: check if layout is a valid matrix

        try:
            with warnings.catch_warnings():
                warnings.filterwarnings("error", category=VisibleDeprecationWarning, message="Creating an ndarray from ragged nested*")
                layout_matrix = np.array(layout_matrix)
        except VisibleDeprecationWarning:
            raise ValueError("The given layout matrix is not valid. Please make sure that all rows have the same length.")

        if len(layout_matrix.shape) == 1:
            layout_matrix = layout_matrix.reshape((1, len(layout_matrix)))  # convert to 2D array

        return layout_matrix

    # -------------------  Fill slide with content  ------------------- #
    def _fill_slide(self):
        """ Fill the slide with content from the internal variables """

        self.set_title()
        self.add_notes()

        # Fill boxes with content
        if len(self.content) > 0:
            self.logger.debug(f"Filling slide with content: {self.content}")
            self.set_layout_matrix()
            self.create_boxes()       # Create boxes based on layout
            self.fill_boxes()         # Fill boxes with content

        # Remove empty placeholders
        if self.remove_placeholders:
            self.remove_empty_ph()

    def set_title(self):
        """ Set the title of the slide. Requires self.title to be set. """

        if self.title is not None:

            # Make sure that title is a string
            try:
                self.title = str(self.title)
            except Exception:
                raise ValueError(f"Could not convert 'title' to a string. The given value is: '{self.title}'.")

            if self._slide.shapes.title is None:
                self.logger.warning("Could not set title of slide. The slide does not have a title box.")
            else:
                self._slide.shapes.title.text = self.title

    def add_notes(self):
        """ Add notes to the slide. """

        if self.notes is not None:

            # Convert notes to a list to enable looping
            if not isinstance(self.notes, list):
                self.notes = [self.notes]

            notes_string = ''
            for s in self.notes:
                if isinstance(s, str):
                    if os.path.exists(s):
                        with open(s, 'r') as f:
                            notes_string += f'\n{f.read()}'
                    else:
                        notes_string += f'\n{s}'
                else:
                    raise ValueError("Notes must be either a string or a list of strings.")

            notes_string = notes_string.lstrip()  # remove leading newline
            self._slide.notes_slide.notes_text_frame.text = notes_string

    def create_boxes(self):
        """ Create boxes for the slide dependent on the internal layout matrix. """

        layout_matrix = self._layout_matrix
        nrows, ncols = layout_matrix.shape

        # Check that margins are valid
        self._validate_margins()

        # Establish left/right/top/bottom margins (in cm)
        left_margin = self.outer_margin if self.left_margin is None else self.left_margin
        right_margin = self.outer_margin if self.right_margin is None else self.right_margin
        top_margin = self.outer_margin if self.top_margin is None else self.top_margin
        bottom_margin = self.outer_margin if self.bottom_margin is None else self.bottom_margin

        # Convert margins from cm to pptx units
        left_margin_unit = Cm(left_margin)
        right_margin_unit = Cm(right_margin)
        top_margin_unit = Cm(top_margin)
        bottom_margin_unit = Cm(bottom_margin)
        inner_margin_unit = Cm(self.inner_margin)

        # Add to top margin based on size of title
        if self._slide.shapes.title.text != "":
            top_margin_unit = self._slide.shapes.title.top + self._slide.shapes.title.height + top_margin_unit

        # How many columns and rows are there?
        n_rows, n_cols = layout_matrix.shape

        # Get total height and width of pictures
        total_width = self._slide_width - left_margin_unit - right_margin_unit - (n_cols - 1) * inner_margin_unit
        total_height = self._slide_height - top_margin_unit - bottom_margin_unit - (n_rows - 1) * inner_margin_unit

        # Check if total_width < 0

        # Get column widths and row heights
        self._validate_ratios()
        if self.width_ratios is None:
            widths = (np.ones(ncols) / ncols) * total_width
        else:
            widths = np.array(self.width_ratios) / sum(self.width_ratios) * total_width

        if self.height_ratios is None:
            heights = (np.ones(nrows) / nrows) * total_height
        else:
            heights = np.array(self.height_ratios) / sum(self.height_ratios) * total_height

        # Box coordinates
        box_numbers = layout_matrix[~np.isnan(layout_matrix)].flatten()
        box_numbers = sorted(set(box_numbers))  # unique box numbers
        box_numbers = [box_number for box_number in box_numbers if box_number >= 0]  # remove negative numbers
        for i in box_numbers:

            # Get column and row number
            coordinates = np.argwhere(layout_matrix == i)

            # Get upper left corner of box
            row, col = coordinates[0]
            left = left_margin_unit + np.sum(widths[:col]) + col * inner_margin_unit
            top = top_margin_unit + np.sum(heights[:row]) + row * inner_margin_unit

            # Get total width and height of box (can span multiple columns and rows)
            width = 0
            height = 0
            rows = set(coordinates[:, 0])
            for row in rows:
                height += heights[row]
            height += (len(rows) - 1) * inner_margin_unit  # add inner margins between rows

            cols = set(coordinates[:, 1])
            for col in cols:
                width += widths[col]
            width += (len(cols) - 1) * inner_margin_unit  # add inner margins between columns

            #  Create box
            self.add_box((left, top, width, height))

    def add_box(self, coordinates):
        """
        Add a box to the slide.

        Parameters
        ----------
        coordinates : tuple
            Coordinates containing (left, top, width, height) of the box (in pptx units).
        """

        box = Box(self._slide, coordinates)
        box.logger = self.logger  # share logger with box

        # Add specific parameters to box
        keys = ["content_alignment", "show_filename", "filename_alignment", "filename_path", "fontsize"]
        parameters = {key: getattr(self, key) for key in keys}
        box.add_parameters(parameters)

        # Add box object to list
        self._boxes.append(box)

    def fill_boxes(self):
        """ Fill the boxes with the elements in self.content """

        for i, element in enumerate(self.content):
            self._boxes[i].fill(element, box_index=i)

    def remove_empty_ph(self):
        """ Remove empty placeholders from the slide. """

        if hasattr(self._slide, 'placeholders'):
            for shape in self._slide.placeholders:
                if shape.has_text_frame and shape.text == '':
                    sp = shape.element
                    sp.getparent().remove(sp)
                    self.logger.debug(f"Removed empty placeholder '{shape.name}' (idx: {shape.placeholder_format.idx})")
